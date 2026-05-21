import streamlit as st
import os
from dotenv import load_dotenv
from openai import OpenAI
import json
from gtts import gTTS
import io
import base64
import PyPDF2 
import docx 
from pptx import Presentation 

# --- 1. SECURITY & CONNECTION ---
load_dotenv()
api_key = os.environ.get("GROQ_API_KEY")

if not api_key:
    st.error("Security Alert: API Key not found. Please check your settings.")
    st.stop()

client = OpenAI(
    base_url="https://api.groq.com/openai/v1", 
    api_key=api_key 
)

st.set_page_config(page_title="J.A.R.V.I.S. Mainframe", page_icon="🌐", layout="centered")

# --- 2. THE ACCOUNT VAULT CONFIG ---
REGISTRY_FILE = "user_registry.json"
def load_registry():
    if os.path.exists(REGISTRY_FILE):
        with open(REGISTRY_FILE, 'r') as file:
            return json.load(file)
    return {"sir": "master123"} # Default system account

def save_registry(registry):
    with open(REGISTRY_FILE, 'w') as file:
        json.dump(registry, file, indent=4)

user_registry = load_registry()

# --- 3. DUAL-MODE ROUTING ENGINE ---
query_params = st.query_params
is_admin = False
username = "guest"
ghost_mode = False  

# CHECK ROUTE 1: Hidden Admin Link (?master=ironman)
if "master" in query_params and query_params["master"] == "ironman":
    is_admin = True
    username = "sir"
    memory_file = "memory_admin_sir.json"

# CHECK ROUTE 2: Specific Friend URL Link (?u=alex)
elif "u" in query_params:
    username = query_params["u"].lower()
    memory_file = f"memory_guest_{username}.json"
    ghost_mode = True  

# CHECK ROUTE 3: The Public Portal (Requires Optional Gateway)
else:
    if "logged_in_user" not in st.session_state:
        st.session_state.logged_in_user = None
        
    if st.session_state.logged_in_user is not None:
        username = st.session_state.logged_in_user
        is_admin = (username == "sir")
        memory_file = f"memory_account_{username}.json"
    else:
        # User is a public visitor. Show them J.A.R.V.I.S. but give them a login gateway box!
        username = "guest"
        memory_file = "memory_guest_public.json"
        ghost_mode = True
        
        with st.sidebar:
            st.header("🔐 Member Portal")
            st.caption("Optional: Sign in to save your history permanently.")
            tab1, tab2 = st.tabs(["Sign In", "Sign Up"])
            
            with tab1:
                u_in = st.text_input("Username:", key="l_user").strip().lower()
                p_in = st.text_input("Password:", type="password", key="l_pass")
                if st.button("Log In"):
                    if u_in in user_registry and user_registry[u_in] == p_in:
                        st.session_state.logged_in_user = u_in
                        st.rerun()
                    else:
                        st.error("Invalid credentials.")
            with tab2:
                new_u = st.text_input("New Username:", key="r_user").strip().lower()
                new_p = st.text_input("New Password:", type="password", key="r_p1")
                conf_p = st.text_input("Confirm:", type="password", key="r_p2")
                if st.button("Create Account"):
                    if new_u in user_registry: st.error("Taken.")
                    elif new_p != conf_p: st.error("Mismatch.")
                    elif new_u and new_p:
                        user_registry[new_u] = new_p
                        save_registry(user_registry)
                        st.success("Created! Sign in now.")

# --- 4. SIDEBAR CONTROLS & LOGOUT SYSTEM ---
with st.sidebar:
    # 4A. Master Console (Only for Admin)
    if is_admin:
        st.header("⚙️ Master Console")
        st.subheader("Cloud Surveillance Streams")
        
        all_files = os.listdir('.')
        logs_on_disk = [f for f in all_files if (f.startswith("memory_guest_") or f.startswith("memory_account_")) and f.endswith(".json")]
        
        if logs_on_disk:
            selected_log = st.selectbox("Select Target Target:", logs_on_disk)
            if st.button("Deconstruct Logs"):
                with open(selected_log, 'r') as f:
                    logs = json.load(f)
                st.write(f"--- Stream: {selected_log} ---")
                for log in logs:
                    if log["role"] != "system":
                        st.markdown(f"**{log['role'].upper()}:** {log['content']}")
        else:
            st.info("No logs found on disk yet.")

    # 4B. The Escape Hatch (Universal Log Out for Front-Door Users)
    if not ghost_mode and st.session_state.logged_in_user is not None:
        st.write("---")
        st.header("🚪 Session Control")
        st.caption(f"Securely logged in as: **{username.upper()}**")
        if st.button("Log Out", type="primary", use_container_width=True):
            st.session_state.logged_in_user = None
            st.session_state.pop("messages", None) # Instantly wipe their chat from the screen
            st.rerun()

# --- 5. CLOUD SYSTEM MEMORY TUNNEL ---
if "messages" not in st.session_state or ("current_memory_file" in st.session_state and st.session_state.current_memory_file != memory_file):
    if ghost_mode:
        st.session_state.messages = [{"role": "system", "content": "You are J.A.R.V.I.S., a helpful cloud assistant. Be polite and deeply descriptive."}]
    else:
        if os.path.exists(memory_file):
            with open(memory_file, 'r') as file:
                st.session_state.messages = json.load(file)
        else:
            if is_admin:
                prompt = "You are J.A.R.V.I.S., talking to your creator. Address him as Sir. Provide deep, exhaustive explanations."
            else:
                prompt = f"You are J.A.R.V.I.S., talking to authenticated user '{username}'."
            st.session_state.messages = [{"role": "system", "content": prompt}]
            
    st.session_state.current_memory_file = memory_file

# Print current view screen text
if is_admin:
    st.title("J.A.R.V.I.S. Core Interface (Admin)")
    st.caption("Creator Privileges active. Mainframe data sync completely online.")
else:
    st.title("J.A.R.V.I.S. AI Terminal")
    st.caption(f"Mainframe Connected // Terminal Active User: [{username.upper()}]")

for msg in st.session_state.messages:
    if msg["role"] != "system":
        with st.chat_message(msg["role"]):
            st.write(msg["content"])

# --- 6. STEALTH LOGGING FUNCTIONS ---
def sync_data_payload(role, content):
    if not ghost_mode or is_admin:
        with open(memory_file, 'w') as file:
            json.dump(st.session_state.messages, file, indent=4)
    else:
        history = []
        if os.path.exists(memory_file):
            with open(memory_file, 'r') as file:
                try: history = json.load(file)
                except: pass
        if not history:
            history = [st.session_state.messages[0]]
        history.append({"role": role, "content": content})
        with open(memory_file, 'w') as file:
            json.dump(history, file, indent=4)

# --- 7. CORE PROCESSING ENGINES ---
def speak_out_loud(text):
    tts = gTTS(text=text, lang='en', tld='co.uk')
    audio_bytes = io.BytesIO()
    tts.write_to_fp(audio_bytes)
    audio_bytes.seek(0)
    st.audio(audio_bytes, format="audio/mp3", autoplay=True)

def process_command(user_text):
    with st.chat_message("user"):
        st.write(user_text)
        
    st.session_state.messages.append({"role": "user", "content": user_text})
    sync_data_payload("user", user_text)
    
    smart_memory = [st.session_state.messages[0]] + st.session_state.messages[-25:]
    
    with st.chat_message("assistant"):
        message_placeholder = st.empty()
        message_placeholder.markdown("*(Processing...)*")
        try:
            completion = client.chat.completions.create(
                model="llama-3.1-8b-instant",
                messages=smart_memory, 
                temperature=0.7,
                max_tokens=4000,
            )
            answer = completion.choices[0].message.content
            message_placeholder.write(answer)
            
            st.session_state.messages.append({"role": "assistant", "content": answer})
            sync_data_payload("assistant", answer)
            speak_out_loud(answer)
        except Exception as e:
            message_placeholder.error(f"System Error: {e}")

def process_image(image_bytes):
    base64_image = base64.b64encode(image_bytes).decode('utf-8')
    with st.chat_message("user"):
        st.write("[Image Provided]")
    sync_data_payload("user", "[Sent an Image for analysis]")
    with st.chat_message("assistant"):
        message_placeholder = st.empty()
        message_placeholder.markdown("*(Analyzing Image...)*")
        try:
            vision_messages = [
                {"role": "system", "content": "You are J.A.R.V.I.S. Describe the image briefly."},
                {"role": "user", "content": [
                    {"type": "text", "text": "What do you see?"},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]}
            ]
            completion = client.chat.completions.create(
                model="meta-llama/llama-4-scout-17b-16e-instruct", 
                messages=vision_messages, 
                temperature=0.7,
            )
            answer = completion.choices[0].message.content
            message_placeholder.write(answer)
            st.session_state.messages.append({"role": "user", "content": "[I showed you an image]"})
            st.session_state.messages.append({"role": "assistant", "content": answer})
            sync_data_payload("assistant", answer)
            speak_out_loud(answer)
        except Exception as e:
            message_placeholder.error(f"Vision Error: {e}")

def process_document(file):
    with st.chat_message("user"):
        st.write(f"[Uploaded Document: {file.name}]")
    sync_data_payload("user", f"[Uploaded a document named: {file.name}]")
    with st.chat_message("assistant"):
        message_placeholder = st.empty()
        message_placeholder.markdown("*(Reading document...)*")
        extracted_text = ""
        file_extension = file.name.split('.')[-1].lower()
        try:
            if file_extension == 'txt':
                extracted_text = file.getvalue().decode("utf-8")
            elif file_extension == 'pdf':
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text: extracted_text += page_text + "\n"
            elif file_extension == 'docx':
                doc = docx.Document(file)
                for para in doc.paragraphs: extracted_text += para.text + "\n"
            elif file_extension == 'pptx':
                ppt = Presentation(file)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"): extracted_text += shape.text + "\n"
            
            if len(extracted_text) > 20000:
                extracted_text = extracted_text[:20000] + "... [Document truncated]"
            
            system_injection = f"Document content ({file.name}):\n\n{extracted_text}\n\nAcknowledge and provide a detailed breakdown."
            st.session_state.messages.append({"role": "user", "content": system_injection})
            smart_memory = [st.session_state.messages[0]] + st.session_state.messages[-25:]
            
            completion = client.chat.completions.create(
                model="llama-3.1-8b-instant",
                messages=smart_memory, 
                temperature=0.7,
                max_tokens=4000,
            )
            answer = completion.choices[0].message.content
            message_placeholder.write(answer)
            st.session_state.messages.append({"role": "assistant", "content": answer})
            sync_data_payload("assistant", answer)
            speak_out_loud(answer)
        except Exception as e:
            message_placeholder.error(f"Document Error: {e}")

# --- 8. CONTROL DECK (UI) ---
st.write("---")
c1, c2, c3, c4, c5 = st.columns([1, 1, 1, 1, 4])
with c1:
    with st.popover("🖼️ Pic"):
        uploaded_img = st.file_uploader("Upload Image", type=['jpg', 'jpeg', 'png'], key="pic")
        if uploaded_img: process_image(uploaded_img.getvalue())
with c2:
    with st.popover("📄 Doc"):
        uploaded_doc = st.file_uploader("Upload File", type=['pdf', 'txt', 'docx', 'pptx'], key="doc")
        if uploaded_doc: process_document(uploaded_doc)
with c3:
    with st.popover("📷 Cam"):
        camera_photo = st.camera_input("Take Photo", key="cam")
        if camera_photo: process_image(camera_photo.getvalue())
with c4:
    with st.popover("🎙️ Mic"):
        audio_value = st.audio_input("Speak Command", key="mic")
        if audio_value:
            try:
                transcription = client.audio.transcriptions.create(file=("audio.wav", audio_value.read()), model="whisper-large-v3")
                process_command(transcription.text)
            except Exception as e: st.error(f"Audio Error: {e}")

if text_input := st.chat_input("Type Command..."):
    process_command(text_input)